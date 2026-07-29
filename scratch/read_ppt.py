import sys
import os

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation

pptx_path = r"C:\Users\HYPE R Series\Downloads\Desain UI Pembelanjaan.pptx"

if not os.path.exists(pptx_path):
    print("File not found:", pptx_path)
    sys.exit(1)

try:
    prs = Presentation(pptx_path)
    text_runs = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        print(f"--- Slide {slide_number} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text.encode('utf-8', 'replace').decode('utf-8'))
                
except Exception as e:
    print("Error reading pptx:", e)
